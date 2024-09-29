import os
import shutil
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from datetime import datetime
from tkinter import Tk
from tkinter.filedialog import askopenfilename

# Ouvrir une boîte de dialogue pour sélectionner le fichier CSV
Tk().withdraw()  # Cacher la fenêtre Tkinter principale
file_path = askopenfilename(title="Sélectionner le fichier CSV", filetypes=[("CSV files", "*.csv")])

# Charger les données du fichier CSV
df = pd.read_csv(file_path, sep=';')

# Obtenir la date actuelle et formater le nom du répertoire
current_date = datetime.now().strftime("%Y-%m-%d")
directory_name = f"{current_date} - Thursday's Meeting"

# Créer le répertoire s'il n'existe pas déjà
if not os.path.exists(directory_name):
    os.makedirs(directory_name)

# Copier le fichier CSV sélectionné vers le nouveau répertoire
csv_dest_path = os.path.join(directory_name, os.path.basename(file_path))
shutil.copy(file_path, csv_dest_path)

# Charger la présentation à partir du modèle
template_path = r'C:\Users\natha\Desktop\vs\TEMPLATE.pptx'  # Mettre à jour le chemin si nécessaire
prs = Presentation(template_path)

# Fonction pour ajouter la date sur une diapositive
def add_date_to_slide(slide):
    date_str = datetime.now().strftime("%Y-%m-%d")  # Format de la date
    textbox = slide.shapes.add_textbox(Cm(17.78), Cm(1.27), Cm(5.08), Cm(1.27))  # Utiliser des cm ici
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = date_str
    p.font.size = Pt(12)  # Taille de la police de la date

# Fonction pour ajouter le numéro de page sur une diapositive
def add_slide_number(slide, number):
    textbox = slide.shapes.add_textbox(Cm(33), Cm(17.27), Cm(2.54), Cm(1.27))  # Utiliser des cm pour la position en bas à droite
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = str(number)
    p.font.size = Pt(12)  # Taille de la police du numéro de page
    p.font.color.rgb = RGBColor(0, 0, 0)  # Couleur du texte

# Ajouter une diapositive pour chaque équipe
def add_team_slide(prs, team_name, team_data):
    # Déterminer le titre en fonction du nom de l'équipe
    if team_name == "Mathieu Palu":
        title_text = "PB&D's Trades"
    elif team_name == "Claire Bernard":
        title_text = "Bank's Trades"
    else:
        title_text = f"Trades de l'équipe : {team_name}"
    
    # Ajouter une diapositive avec un titre (utiliser le layout du modèle)
    slide_layout = prs.slide_layouts[5]  # Diapositive vide
    slide = prs.slides.add_slide(slide_layout)
    
    # Ajouter le titre de la diapositive
    title = slide.shapes.title
    title.text = title_text
    
    # Définir la position du tableau sur la diapositive
    x, y, cx, cy = Cm(1.27), Cm(5), Cm(21.59), Cm(12.7)  # Utilisation de cm ici
    
    # Déterminer le nombre de lignes et de colonnes pour le tableau
    rows, cols = len(team_data) + 1, len(team_data.columns)
    
    # Ajouter un tableau
    table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table
    
    # Définir la taille de la police
    font_size = Pt(10)
    green_color = RGBColor(0, 128, 0)
    red_color = RGBColor(255, 0, 0)  # Rouge
    orange_color = RGBColor(255, 165, 0)  # Orange

    # Ajouter les en-têtes de colonnes
    for col_idx, col_name in enumerate(team_data.columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = green_color
        
        # Changer la taille de la police de l'en-tête
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = font_size

    # Ajouter les données du tableau
    for row_idx, trade in enumerate(team_data.values):
        for col_idx, value in enumerate(trade):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            # Changer la taille de la police pour chaque cellule de données
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = font_size

            # Appliquer la couleur de fond de la cellule selon la valeur
            if str(value).lower() == "high":
                cell.fill.solid()
                cell.fill.fore_color.rgb = red_color  # Rouge pour "high"
            elif str(value).lower() == "medium":
                cell.fill.solid()
                cell.fill.fore_color.rgb = orange_color  # Orange pour "medium"
    
    # Ajouter l'image si l'équipe est "Julien Dupont"
    if team_name == "Julien Dupont":
        img_path = r"C:\Users\natha\Desktop\vs\iamge1.jpg"  # Mettre à jour le chemin si nécessaire
        slide.shapes.add_picture(img_path, Cm(20.32), Cm(10), width=Cm(20), height=Cm(2.54))  # Utiliser des cm ici
    
    # Ajouter le texte "5 trades" pour Sophie Martin
    if team_name == "Sophie Martin":
        textbox = slide.shapes.add_textbox(Cm(25.4), Cm(1.27), Cm(5.08), Cm(1.27))  # Position en haut à droite en cm
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = "5 trades"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.font.name = 'Arial'

    return slide

# Liste des équipes uniques
teams = df['Equipe'].unique()
print(teams)

# Ajouter une diapositive pour chaque équipe
slides_to_add = []

for team in teams:
    team_data = df[df['Equipe'] == team]
    slide = add_team_slide(prs, team, team_data)
    slides_to_add.append(slide)

# Réorganiser les slides pour les insérer entre la première et la deuxième
for slide in slides_to_add:
    prs.slides._sldIdLst.insert(1, prs.slides._sldIdLst[-1])

# Ajouter la date à la première diapositive
add_date_to_slide(prs.slides[0])

# Ajouter le numéro de page à chaque diapositive
for i, slide in enumerate(prs.slides):
    add_slide_number(slide, i + 1)

# Sauvegarder la présentation PowerPoint dans le même répertoire que le fichier CSV
output_path = os.path.join(directory_name, 'pipeline_presentation.pptx')
prs.save(output_path)

print(f"Présentation PowerPoint et fichier CSV enregistrés dans le répertoire : {directory_name}")
